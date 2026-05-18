from io import BytesIO

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

from app.core.exceptions import ParsingError


def extract_text(file_bytes: bytes, file_type: str) -> str:
    try:
        if file_type == "pdf":
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            return "\n\n".join(page.get_text() for page in doc)

        if file_type == "pptx":
            prs = Presentation(BytesIO(file_bytes))
            return "\n\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if shape.has_text_frame
            )

        if file_type == "docx":
            doc = DocxDocument(BytesIO(file_bytes))
            return "\n\n".join(p.text for p in doc.paragraphs)

    except Exception as e:
        raise ParsingError(f"Failed to extract text from {file_type}: {e}") from e

    raise ParsingError(f"Unsupported file type: {file_type}")


def count_pages(file_bytes: bytes, file_type: str) -> int:
    try:
        if file_type == "pdf":
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            return doc.page_count
        if file_type == "pptx":
            prs = Presentation(BytesIO(file_bytes))
            return len(prs.slides)
        if file_type == "docx":
            return count_docx_words(file_bytes)
    except Exception as e:
        raise ParsingError(f"Failed to count pages in {file_type}: {e}") from e
    raise ParsingError(f"Unsupported file type: {file_type}")


def count_docx_words(file_bytes: bytes) -> int:
    try:
        doc = DocxDocument(BytesIO(file_bytes))
        words = [word for p in doc.paragraphs for word in p.text.split() if word.strip()]
        return len(words)
    except Exception as e:
        raise ParsingError(f"Failed to count words in docx: {e}") from e
